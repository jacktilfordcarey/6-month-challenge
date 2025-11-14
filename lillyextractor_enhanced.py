"""
Lilly Data Intelligence Hub - Enhanced Version
Supports Excel, CSV, PowerPoint, Text, JSON with AI chatbot and AR visualizations
"""

import streamlit as st
import pandas as pd
import os
import google.generativeai as genai
from dotenv import load_dotenv
import plotly.graph_objects as go
from io import BytesIO
import json

# Try to import PowerPoint support
try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False
    st.warning("⚠️ PowerPoint support not installed. Run: pip install python-pptx")

# Load environment variables
load_dotenv()

# Configure page
st.set_page_config(
    page_title="Lilly Data Intelligence Hub",
    page_icon="💊",
    layout="wide"
)

# Configure Gemini
GEMINI_API_KEY = os.getenv('GEMINI_API_KEY') or os.getenv('GOOGLE_API_KEY')
if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)
    model = genai.GenerativeModel('gemini-2.0-flash-exp')
else:
    model = None

# Custom CSS for Lilly theme
st.markdown("""
<style>
    .main {
        background-color: #1a1a1a !important;
        padding: 2rem;
    }
    [data-testid="stSidebar"] {
        background-color: #0d0d0d !important;
    }
    .stApp {
        background-color: #0d0d0d !important;
    }
    header[data-testid="stHeader"] {
        background-color: #0d0d0d !important;
    }
    h1, h2, h3 {
        color: #E41E26 !important;
    }
    .stButton>button {
        background-color: #E41E26 !important;
        color: white !important;
        border: none !important;
        border-radius: 5px;
        padding: 0.5rem 1rem;
        font-weight: bold;
    }
    .stButton>button:hover {
        background-color: #c91821 !important;
    }
    .stMetric {
        background-color: #2d2d2d;
        padding: 1rem;
        border-radius: 8px;
        border-left: 4px solid #E41E26;
    }
    .stExpander {
        background-color: #2d2d2d !important;
        border: 1px solid #3d3d3d !important;
        border-radius: 8px;
    }
    .stChatMessage {
        background-color: #2d2d2d !important;
        border-radius: 8px;
    }
</style>
""", unsafe_allow_html=True)

# Header
st.markdown("""
<div style='text-align: center; padding: 2rem 0;'>
    <h1 style='color: #E41E26; font-size: 3rem; margin-bottom: 0.5rem;'>💊 Lilly Data Intelligence Hub</h1>
    <p style='color: #4a4a4a; font-size: 1.2rem;'>Upload data, get insights, visualizations, and AR experiences</p>
</div>
""", unsafe_allow_html=True)

# Helper functions
def extract_text_from_pptx(file):
    """Extract text and tables from PowerPoint"""
    if not PPTX_SUPPORT:
        return "PowerPoint support not available. Install python-pptx.", []
    
    prs = Presentation(file)
    text_content = []
    tables_data = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
            
            if shape.has_table:
                table = shape.table
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                if table_data:
                    tables_data.append({"slide": slide_num, "data": table_data})
        
        if slide_text:
            text_content.append(f"Slide {slide_num}:\n" + "\n".join(slide_text))
    
    return "\n\n".join(text_content), tables_data

def parse_json_to_dataframe(file):
    """Parse JSON file to DataFrame"""
    try:
        content = json.load(file)
        if isinstance(content, list):
            return pd.DataFrame(content)
        elif isinstance(content, dict):
            for key, value in content.items():
                if isinstance(value, list) and len(value) > 0:
                    return pd.DataFrame(value)
            return pd.DataFrame([content])
        return pd.DataFrame()
    except Exception as e:
        st.error(f"Error parsing JSON: {str(e)}")
        return pd.DataFrame()

# Initialize session state
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []
if 'analysis_text' not in st.session_state:
    st.session_state.analysis_text = ''

# File upload section
st.markdown("### 📁 Upload Your Data")

file_type = st.radio(
    "Select input type:",
    ["Data File (Excel/CSV)", "PowerPoint Presentation", "Text/Analysis Document", "JSON Data"],
    horizontal=True
)

uploaded_file = None

if file_type == "Data File (Excel/CSV)":
    uploaded_file = st.file_uploader(
        "Choose an Excel or CSV file",
        type=['xlsx', 'xls', 'csv'],
        help="Upload your data file for analysis"
    )
elif file_type == "PowerPoint Presentation":
    uploaded_file = st.file_uploader(
        "Choose a PowerPoint file",
        type=['pptx'],
        help="Upload PowerPoint with data tables or analysis"
    )
elif file_type == "Text/Analysis Document":
    uploaded_file = st.file_uploader(
        "Choose a text file",
        type=['txt', 'md'],
        help="Upload text document with analysis or data"
    )
else:  # JSON Data
    uploaded_file = st.file_uploader(
        "Choose a JSON file",
        type=['json'],
        help="Upload JSON file with structured data"
    )

# Main processing
if uploaded_file is not None:
    try:
        df = None
        extracted_text = None
        
        # Load data based on file type
        if file_type == "Data File (Excel/CSV)":
            if uploaded_file.name.endswith('.csv'):
                df = pd.read_csv(uploaded_file)
            else:
                df = pd.read_excel(uploaded_file)
            st.success(f"✅ Loaded data: {len(df)} rows × {len(df.columns)} columns")
        
        elif file_type == "PowerPoint Presentation":
            extracted_text, tables = extract_text_from_pptx(uploaded_file)
            st.markdown("### 📊 Extracted Content")
            with st.expander("View extracted text", expanded=False):
                st.text_area("PowerPoint Content", extracted_text, height=300, key="pptx_text")
            
            if tables:
                st.markdown("### 📋 Extracted Tables")
                for idx, table_info in enumerate(tables):
                    st.markdown(f"**Table from Slide {table_info['slide']}**")
                    table_data = table_info['data']
                    if len(table_data) > 1:
                        df_temp = pd.DataFrame(table_data[1:], columns=table_data[0])
                        st.dataframe(df_temp, use_container_width=True)
                        if df is None:
                            df = df_temp
        
        elif file_type == "Text/Analysis Document":
            extracted_text = uploaded_file.read().decode('utf-8')
            st.markdown("### 📄 Document Content")
            with st.expander("View document text", expanded=False):
                st.text_area("Text Content", extracted_text, height=300, key="text_content")
        
        elif file_type == "JSON Data":
            df = parse_json_to_dataframe(uploaded_file)
            if not df.empty:
                st.success(f"✅ Loaded JSON data: {len(df)} records")
        
        # Display preview
        st.markdown("### 📊 Data Preview")
        
        if df is not None and not df.empty:
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("Rows", f"{len(df):,}")
            with col2:
                st.metric("Columns", len(df.columns))
            with col3:
                st.metric("Memory", f"{df.memory_usage(deep=True).sum() / 1024:.1f} KB")
            
            with st.expander("View Data Sample", expanded=True):
                st.dataframe(df.head(10), use_container_width=True)
        
        elif extracted_text:
            col1, col2 = st.columns(2)
            with col1:
                st.metric("Characters", f"{len(extracted_text):,}")
            with col2:
                st.metric("Words", f"{len(extracted_text.split()):,}")
        
        # AI Analysis
        st.markdown("### 🤖 AI Analysis")
        
        if st.button("Generate AI Insights", type="primary", use_container_width=True):
            if not model:
                st.error("⚠️ AI model not configured. Please set GEMINI_API_KEY in .env file")
            else:
                with st.spinner("Analyzing with AI..."):
                    if df is not None and not df.empty:
                        data_summary = f"""
                        Dataset Overview:
                        - Rows: {len(df)}
                        - Columns: {len(df.columns)}
                        - Column Names: {', '.join(df.columns.tolist())}
                        
                        Data Types:
                        {df.dtypes.to_string()}
                        
                        Statistical Summary:
                        {df.describe().to_string()}
                        
                        First few rows:
                        {df.head().to_string()}
                        """
                    elif extracted_text:
                        data_summary = f"""
                        Content Analysis:
                        - Type: {file_type}
                        - Length: {len(extracted_text)} characters
                        
                        Content:
                        {extracted_text[:2000]}
                        """
                    else:
                        data_summary = "No data available."
                    
                    prompt = f"""Analyze this data and provide key insights, trends, and recommendations.
                    Focus on actionable insights for healthcare/pharmaceutical context.
                    
                    {data_summary}
                    """
                    
                    try:
                        response = model.generate_content(prompt)
                        analysis_text = response.text
                        st.session_state.analysis_text = analysis_text
                        
                        st.markdown("#### 📈 Analysis Results")
                        st.markdown(analysis_text)
                    except Exception as e:
                        st.error(f"Error generating analysis: {str(e)}")
        
        # Show stored analysis
        if st.session_state.analysis_text:
            with st.expander("View Previous Analysis", expanded=False):
                st.markdown(st.session_state.analysis_text)
        
        # Visualizations
        if df is not None and not df.empty:
            if st.button("Generate Visualizations", type="primary", use_container_width=True):
                with st.spinner("Creating visualizations..."):
                    st.markdown("### 📊 Data Visualizations")
                    
                    numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
                    categorical_cols = df.select_dtypes(include=['object', 'category']).columns.tolist()
                    
                    # 3D Scatter Plot
                    if len(numeric_cols) >= 2:
                        st.markdown("#### 🎯 3D Interactive Scatter Plot")
                        
                        x_col = numeric_cols[0]
                        y_col = numeric_cols[1] if len(numeric_cols) > 1 else numeric_cols[0]
                        z_col = numeric_cols[2] if len(numeric_cols) > 2 else numeric_cols[0]
                        
                        color_col = categorical_cols[0] if categorical_cols else None
                        
                        if color_col:
                            unique_categories = df[color_col].unique()
                            colors = ['#E41E26' if i % 2 == 0 else '#4a4a4a' for i in range(len(unique_categories))]
                            color_map = {cat: colors[i] for i, cat in enumerate(unique_categories)}
                            
                            fig = go.Figure(data=[go.Scatter3d(
                                x=df[x_col],
                                y=df[y_col],
                                z=df[z_col],
                                mode='markers',
                                marker=dict(
                                    size=6,
                                    color=[color_map[val] for val in df[color_col]],
                                    line=dict(color='white', width=0.5)
                                ),
                                text=[f"{color_col}: {val}" for val in df[color_col]],
                                hovertemplate=f'<b>%{{text}}</b><br>{x_col}: %{{x}}<br>{y_col}: %{{y}}<br>{z_col}: %{{z}}<extra></extra>'
                            )])
                        else:
                            fig = go.Figure(data=[go.Scatter3d(
                                x=df[x_col],
                                y=df[y_col],
                                z=df[z_col],
                                mode='markers',
                                marker=dict(
                                    size=6,
                                    color='#E41E26',
                                    line=dict(color='white', width=0.5)
                                ),
                                hovertemplate=f'{x_col}: %{{x}}<br>{y_col}: %{{y}}<br>{z_col}: %{{z}}<extra></extra>'
                            )])
                        
                        fig.update_layout(
                            title=dict(text=f'3D Scatter: {x_col} vs {y_col} vs {z_col}', font=dict(color='#E41E26', size=18)),
                            scene=dict(
                                xaxis=dict(title=x_col, gridcolor='#3d3d3d', backgroundcolor='#1a1a1a'),
                                yaxis=dict(title=y_col, gridcolor='#3d3d3d', backgroundcolor='#1a1a1a'),
                                zaxis=dict(title=z_col, gridcolor='#3d3d3d', backgroundcolor='#1a1a1a'),
                            ),
                            paper_bgcolor='#1a1a1a',
                            font=dict(color='#ffffff'),
                            height=600
                        )
                        st.plotly_chart(fig, use_container_width=True)
                    
                    # Distribution plots
                    if numeric_cols:
                        st.markdown("#### 📈 Distribution Analysis")
                        for col in numeric_cols[:3]:
                            fig = go.Figure(data=[go.Histogram(
                                x=df[col],
                                marker_color='#E41E26',
                                marker_line_color='white',
                                marker_line_width=1.5
                            )])
                            fig.update_layout(
                                title=dict(text=f'Distribution of {col}', font=dict(color='#E41E26', size=16)),
                                xaxis=dict(title=col, gridcolor='#3d3d3d'),
                                yaxis=dict(title='Frequency', gridcolor='#3d3d3d'),
                                paper_bgcolor='#1a1a1a',
                                plot_bgcolor='#1a1a1a',
                                font=dict(color='#ffffff'),
                                height=400
                            )
                            st.plotly_chart(fig, use_container_width=True)
                    
                    # Correlation heatmap
                    if len(numeric_cols) > 1:
                        st.markdown("#### 🔥 Correlation Matrix")
                        corr_matrix = df[numeric_cols].corr()
                        fig = go.Figure(data=go.Heatmap(
                            z=corr_matrix.values,
                            x=corr_matrix.columns,
                            y=corr_matrix.columns,
                            colorscale=[[0, '#000000'], [0.5, '#ffffff'], [1, '#E41E26']],
                            zmid=0,
                            text=corr_matrix.values.round(2),
                            texttemplate='%{text}',
                            textfont={"size": 10}
                        ))
                        fig.update_layout(
                            title=dict(text='Feature Correlation Heatmap', font=dict(color='#E41E26', size=18)),
                            paper_bgcolor='#1a1a1a',
                            plot_bgcolor='#1a1a1a',
                            font=dict(color='#ffffff'),
                            height=600
                        )
                        st.plotly_chart(fig, use_container_width=True)
        
        # AR Visualization Export
        st.markdown("---")
        st.markdown("### 🥽 AR Visualization Export")
        
        if df is not None and not df.empty:
            numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
            if len(numeric_cols) >= 2:
                if st.button("Generate AR-Ready 3D Model", use_container_width=True):
                    with st.spinner("Creating AR visualization..."):
                        x_col = numeric_cols[0]
                        y_col = numeric_cols[1] if len(numeric_cols) > 1 else numeric_cols[0]
                        z_col = numeric_cols[2] if len(numeric_cols) > 2 else numeric_cols[0]
                        
                        fig = go.Figure(data=[go.Scatter3d(
                            x=df[x_col], y=df[y_col], z=df[z_col],
                            mode='markers',
                            marker=dict(
                                size=8,
                                color=df[z_col] if z_col in df.columns else '#E41E26',
                                colorscale=[[0, '#000000'], [0.5, '#4a4a4a'], [1, '#E41E26']],
                                line=dict(color='white', width=1),
                                showscale=True
                            )
                        )])
                        
                        fig.update_layout(
                            title=dict(text='AR-Ready 3D Visualization', font=dict(color='#E41E26', size=20)),
                            scene=dict(
                                xaxis=dict(title=x_col, backgroundcolor='#1a1a1a', gridcolor='#3d3d3d'),
                                yaxis=dict(title=y_col, backgroundcolor='#1a1a1a', gridcolor='#3d3d3d'),
                                zaxis=dict(title=z_col, backgroundcolor='#1a1a1a', gridcolor='#3d3d3d'),
                                camera=dict(eye=dict(x=1.5, y=1.5, z=1.5))
                            ),
                            paper_bgcolor='#1a1a1a',
                            font=dict(color='#ffffff'),
                            height=700
                        )
                        
                        st.plotly_chart(fig, use_container_width=True)
                        
                        html_buffer = BytesIO()
                        fig.write_html(html_buffer)
                        html_buffer.seek(0)
                        
                        st.download_button(
                            label="📱 Download AR Visualization (HTML)",
                            data=html_buffer,
                            file_name="ar_visualization.html",
                            mime="text/html"
                        )
                        st.info("💡 **AR Tip:** Open on phone/tablet with AR capabilities for immersive viewing")
        
        # Chatbot
        st.markdown("---")
        st.markdown("### 💬 Ask Questions About Your Data")
        
        for message in st.session_state.chat_history:
            with st.chat_message(message["role"]):
                st.markdown(message["content"])
        
        if user_question := st.chat_input("Ask anything about your data..."):
            st.session_state.chat_history.append({"role": "user", "content": user_question})
            
            with st.chat_message("user"):
                st.markdown(user_question)
            
            with st.chat_message("assistant"):
                if not model:
                    answer = "⚠️ AI not configured. Set GEMINI_API_KEY."
                    st.markdown(answer)
                else:
                    with st.spinner("Thinking..."):
                        if df is not None and not df.empty:
                            context = f"""
                            Dataset: {len(df)} rows, {len(df.columns)} columns
                            Columns: {', '.join(df.columns.tolist())}
                            Stats: {df.describe().to_string()}
                            Previous analysis: {st.session_state.analysis_text}
                            
                            Question: {user_question}
                            """
                        elif extracted_text:
                            context = f"""
                            Content: {extracted_text[:1500]}
                            Previous analysis: {st.session_state.analysis_text}
                            
                            Question: {user_question}
                            """
                        else:
                            context = user_question
                        
                        try:
                            response = model.generate_content(context)
                            answer = response.text
                            st.markdown(answer)
                        except Exception as e:
                            answer = f"Error: {str(e)}"
                            st.error(answer)
                
                st.session_state.chat_history.append({"role": "assistant", "content": answer})
        
        if st.button("Clear Chat"):
            st.session_state.chat_history = []
            st.rerun()
        
        # Downloads
        st.markdown("---")
        st.markdown("### 📥 Download Options")
        
        col1, col2, col3 = st.columns(3)
        
        if df is not None and not df.empty:
            with col1:
                csv = df.to_csv(index=False).encode('utf-8')
                st.download_button("📊 Download Data (CSV)", csv, "data_export.csv", "text/csv")
            
            with col2:
                if len(df.select_dtypes(include=['number']).columns) > 0:
                    stats = df.describe().to_csv().encode('utf-8')
                    st.download_button("📈 Download Statistics", stats, "statistics.csv", "text/csv")
        
        with col3:
            if st.session_state.analysis_text:
                report = f"""Lilly Data Intelligence Hub - Report
=====================================
File: {uploaded_file.name}
Type: {file_type}

AI Analysis:
{st.session_state.analysis_text}

Chat History:
"""
                for msg in st.session_state.chat_history:
                    report += f"\n{msg['role'].upper()}: {msg['content']}\n"
                
                st.download_button("📄 Download Report", report.encode('utf-8'), "report.txt", "text/plain")
    
    except Exception as e:
        st.error(f"Error: {str(e)}")
        st.info("Please ensure your file is valid and properly formatted")

else:
    st.info("👆 Upload a file to get started")
    st.markdown("""
    **Supported formats:**
    - 📊 Data: Excel (.xlsx, .xls), CSV
    - 📽️ PowerPoint (.pptx) - extracts tables and text
    - 📄 Text: .txt, .md files
    - 🔧 JSON structured data
    """)
