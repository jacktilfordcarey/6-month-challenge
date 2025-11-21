"""
MedInsight - AI Tool to Help Understand/Visuals RWE
Clean, professional interface for RWE data analysis and visualization
"""

import streamlit as st
import pandas as pd
import os
import google.generativeai as genai
from dotenv import load_dotenv
import plotly.graph_objects as go
from io import BytesIO
import json

# Try to import AI libraries
try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False

try:
    from PyPDF2 import PdfReader
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False

try:
    from groq import Groq
    GROQ_AVAILABLE = True
except ImportError:
    GROQ_AVAILABLE = False

try:
    from openai import OpenAI
    OPENAI_AVAILABLE = True
except ImportError:
    OPENAI_AVAILABLE = False

# Load environment variables from .env file
load_dotenv()

# Configure page
st.set_page_config(
    page_title="MedInsight",
    page_icon="🔬",
    layout="wide"
)

# Initialize session state variables FIRST
if 'high_contrast' not in st.session_state:
    st.session_state.high_contrast = False
if 'magnifier' not in st.session_state:
    st.session_state.magnifier = False
if 'tts_enabled' not in st.session_state:
    st.session_state.tts_enabled = False
if 'chatbot_messages' not in st.session_state:
    st.session_state.chatbot_messages = []
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []
if 'analysis_text' not in st.session_state:
    st.session_state.analysis_text = ''
if 'input_method' not in st.session_state:
    st.session_state.input_method = "Upload File"

# Configure AI - Only using Groq (keys stored in .env file, never in code)
GROQ_API_KEY = os.getenv('GROQ_API_KEY')
GEMINI_API_KEY = None  # Removed - use .env if needed
OPENAI_API_KEY = None  # Removed - use .env if needed

model = None
active_model = None

# Try Groq first (with SSL verification disabled for corporate networks)
if GROQ_API_KEY and GROQ_AVAILABLE:
    try:
        import httpx
        http_client = httpx.Client(verify=False, timeout=60.0)
        groq_client = Groq(api_key=GROQ_API_KEY, http_client=http_client)
        model = groq_client
        active_model = "groq"
    except:
        pass

# Try Gemini second
if model is None and GEMINI_API_KEY:
    try:
        genai.configure(api_key=GEMINI_API_KEY)
        model = genai.GenerativeModel('gemini-2.0-flash-exp')
        active_model = "gemini"
    except:
        pass

# Try OpenAI third (with SSL verification disabled for corporate networks)
if model is None and OPENAI_API_KEY and OPENAI_AVAILABLE:
    try:
        import httpx
        http_client = httpx.Client(verify=False, timeout=60.0)
        openai_client = OpenAI(api_key=OPENAI_API_KEY, http_client=http_client)
        model = openai_client
        active_model = "openai"
    except:
        pass

def generate_ai_response(prompt):
    """Generate AI response with fallback support"""
    global model, active_model
    
    if model is None:
        return "AI not configured. Please set API keys."
    
    errors_encountered = []
    
    # Try primary model first
    try:
        if active_model == "groq":
            # Use llama-3.1-8b-instant (supported and fast)
            response = model.chat.completions.create(
                model="llama-3.1-8b-instant",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.7,
                max_tokens=2000,
                timeout=60
            )
            return response.choices[0].message.content
        
        elif active_model == "gemini":
            response = model.generate_content(prompt)
            return response.text
        
        elif active_model == "openai":
            response = model.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.7,
                max_tokens=2000,
                timeout=30
            )
            return response.choices[0].message.content
    
    except Exception as e:
        error_msg = str(e).lower()
        errors_encountered.append(f"{active_model}: {str(e)}")
        
        # Try Groq fallback (if not already using it) - try multiple times
        if active_model != "groq" and GROQ_API_KEY and GROQ_AVAILABLE:
            for attempt in range(2):  # Try twice
                try:
                    import httpx
                    http_client = httpx.Client(verify=False, timeout=60.0)
                    groq_client = Groq(api_key=GROQ_API_KEY, http_client=http_client)
                    response = groq_client.chat.completions.create(
                        model="llama-3.1-8b-instant",
                        messages=[{"role": "user", "content": prompt}],
                        temperature=0.7,
                        max_tokens=2000
                    )
                    model = groq_client
                    active_model = "groq"
                    return response.choices[0].message.content
                except Exception as groq_error:
                    if attempt == 1:  # Only log on last attempt
                        errors_encountered.append(f"Groq (retry {attempt+1}): {str(groq_error)}")
                    import time
                    time.sleep(1)  # Wait 1 second between retries
        
        # Try Gemini fallback (if not already using it and not quota error)
        if active_model != "gemini" and GEMINI_API_KEY and "quota" not in error_msg and "429" not in error_msg:
            try:
                genai.configure(api_key=GEMINI_API_KEY)
                fallback_model = genai.GenerativeModel('gemini-2.0-flash-exp')
                response = fallback_model.generate_content(prompt)
                model = fallback_model
                active_model = "gemini"
                return response.text
            except Exception as gemini_error:
                errors_encountered.append(f"Gemini: {str(gemini_error)}")
        
        # Try OpenAI fallback (if not already using it)
        if active_model != "openai" and OPENAI_API_KEY and OPENAI_AVAILABLE:
            try:
                import httpx
                http_client = httpx.Client(verify=False, timeout=60.0)
                openai_client = OpenAI(api_key=OPENAI_API_KEY, http_client=http_client)
                response = openai_client.chat.completions.create(
                    model="gpt-3.5-turbo",
                    messages=[{"role": "user", "content": prompt}],
                    temperature=0.7,
                    max_tokens=2000
                )
                model = openai_client
                active_model = "openai"
                return response.choices[0].message.content
            except Exception as openai_error:
                errors_encountered.append(f"OpenAI: {str(openai_error)}")
        
        # Return all errors for debugging
        return f"All AI services failed. Errors: {' | '.join(errors_encountered)}"

# Load external CSS file
def load_css():
    with open("style.css") as f:
        st.markdown(f'<style>{f.read()}</style>', unsafe_allow_html=True)

try:
    load_css()
except:
    # Fallback inline CSS if file not found
    st.markdown("""
    <style>
        .main {
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%) !important;
        }
        .block-container {
            background: rgba(255, 255, 255, 0.95) !important;
            border-radius: 20px !important;
            padding: 2rem !important;
        }
        [data-testid="stRadio"] input[type="radio"] {
            accent-color: #667eea !important;
        }
    </style>
    """, unsafe_allow_html=True)

# Keep minimal inline CSS for compatibility
st.markdown("""
<style>
    /* Ensure main styling */
    .main {
        padding: 2rem;
        max-width: 100% !important;
    }
    
    /* Better spacing when sidebar is open */
    .block-container {
        padding-left: 2rem !important;
        padding-right: 2rem !important;
        max-width: 100% !important;
    }
    
    /* Sidebar */
    [data-testid="stSidebar"] {
        background-color: #E3F2FD !important;
        min-width: 280px !important;
    }
    
    /* Reduce sidebar width when collapsed */
    [data-testid="stSidebar"][aria-expanded="false"] {
        min-width: 0px !important;
    }
    
    /* App background - White */
    .stApp {
        background-color: #FFFFFF !important;
    }
    
    /* Header - Light Blue */
    header[data-testid="stHeader"] {
        background-color: #BBDEFB !important;
    }
    
    /* Toolbar - Light Blue */
    [data-testid="stToolbar"] {
        background-color: #BBDEFB !important;
    }
    
    /* Hide deploy button in header */
    [data-testid="stToolbar"] button[kind="header"] {
        display: none !important;
    }
    
    /* Top decoration */
    .stApp > header {
        background-color: #BBDEFB !important;
    }
    
    /* Block container - slightly lighter grey */
    [data-testid="stVerticalBlock"] {
        background-color: transparent !important;
    }
    
    /* All text default to dark grey/black */
    body, p, span, div, label, .stMarkdown, .stText {
        color: #212121 !important;
    }
    
    /* Headers - Blue */
    h1 {
        color: #1565C0 !important;
        font-weight: 700;
        margin-bottom: 0.5rem !important;
        word-wrap: break-word !important;
    }
    h2, h3 {
        color: #1976D2 !important;
        font-weight: 600;
        margin-top: 2rem !important;
        word-wrap: break-word !important;
    }
    
    /* Prevent text overflow */
    p, div, span {
        word-wrap: break-word !important;
        overflow-wrap: break-word !important;
    }
    
    /* Subtitle text */
    .stMarkdown h3 {
        color: #424242 !important;
        font-weight: 400;
    }
    
    /* Buttons - All buttons blue */
    button[kind="primary"],
    button[kind="secondary"],
    .stButton button,
    .stDownloadButton button {
        background-color: #2196F3 !important;
        color: #FFFFFF !important;
        border: none !important;
        padding: 0.6rem 1.5rem !important;
        font-size: 1rem !important;
        border-radius: 6px !important;
        font-weight: 600 !important;
        transition: all 0.3s ease !important;
    }
    
    button[kind="primary"]:hover,
    button[kind="secondary"]:hover,
    .stButton button:hover,
    .stDownloadButton button:hover {
        background-color: #1976D2 !important;
        box-shadow: 0 4px 12px rgba(33, 150, 243, 0.4) !important;
        transform: translateY(-1px) !important;
    }
    
    /* File uploader - Light grey box with blue border */
    [data-testid="stFileUploader"] {
        background-color: #F5F5F5 !important;
        border: 2px dashed #2196F3 !important;
        border-radius: 8px !important;
        padding: 2rem !important;
    }
    [data-testid="stFileUploader"] label {
        color: #212121 !important;
        font-size: 1.1rem !important;
    }
    [data-testid="stFileUploader"] section {
        background-color: #F5F5F5 !important;
    }
    [data-testid="stFileUploader"] small {
        color: #757575 !important;
    }
    
    /* Dataframe - Light background */
    [data-testid="stDataFrame"] {
        background-color: transparent !important;
        border-radius: 8px !important;
        padding: 1rem !important;
        border: 1px solid #E0E0E0 !important;
    }
    .stDataFrame {
        background-color: transparent !important;
    }
    
    /* Dataframe table styling */
    [data-testid="stDataFrame"] table {
        background-color: transparent !important;
    }
    
    /* Table header - blue */
    [data-testid="stDataFrame"] thead tr th {
        background-color: #2196F3 !important;
        color: #FFFFFF !important;
        font-weight: 600 !important;
        border-bottom: 2px solid #1976D2 !important;
        padding: 0.75rem !important;
    }
    
    /* Keep headers blue on dropdown/filter */
    [data-testid="stDataFrame"] thead tr th:hover {
        background-color: #1E88E5 !important;
    }
    
    [data-testid="stDataFrame"] thead tr th[aria-sort] {
        background-color: #1E88E5 !important;
    }
    
    /* Table body cells - light grey */
    [data-testid="stDataFrame"] tbody tr td {
        background-color: #FAFAFA !important;
        color: #212121 !important;
        border-bottom: 1px solid #E0E0E0 !important;
        padding: 0.5rem !important;
    }
    
    /* Alternating row colors for better readability */
    [data-testid="stDataFrame"] tbody tr:nth-child(even) td {
        background-color: #F5F5F5 !important;
    }
    
    /* Hover effect on rows */
    [data-testid="stDataFrame"] tbody tr:hover td {
        background-color: #E3F2FD !important;
    }
    
    /* Index column styling */
    [data-testid="stDataFrame"] tbody tr th {
        background-color: #BBDEFB !important;
        color: #212121 !important;
        font-weight: 500 !important;
    }
    
    /* Expander - Light grey with blue accent */
    [data-testid="stExpander"] {
        background-color: #F5F5F5 !important;
        border: 1px solid #E0E0E0 !important;
        border-radius: 8px !important;
        margin: 1rem 0 !important;
    }
    details[open] > summary {
        border-bottom: 2px solid #2196F3 !important;
        padding-bottom: 0.5rem !important;
    }
    .streamlit-expanderHeader {
        background-color: #F5F5F5 !important;
        color: #212121 !important;
        font-weight: 600 !important;
    }
    
    /* Metrics - Light grey cards with blue accent */
    [data-testid="stMetric"] {
        background-color: #F5F5F5 !important;
        padding: 1.2rem !important;
        border-radius: 8px !important;
        border-left: 4px solid #2196F3 !important;
        box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1) !important;
    }
    [data-testid="stMetricValue"] {
        color: #1565C0 !important;
        font-weight: 700 !important;
        font-size: 2rem !important;
    }
    [data-testid="stMetricLabel"] {
        color: #424242 !important;
        font-weight: 500 !important;
    }
    
    /* Info/Success/Error boxes */
    .stAlert {
        background-color: #E3F2FD !important;
        color: #212121 !important;
        border-left: 5px solid #2196F3 !important;
        border-radius: 6px !important;
        padding: 1rem !important;
    }
    
    /* Success message */
    .stSuccess {
        background-color: #E8F5E9 !important;
        color: #212121 !important;
        border-left: 5px solid #4CAF50 !important;
    }
    
    /* Info message */
    .stInfo {
        background-color: #E3F2FD !important;
        color: #212121 !important;
        border-left: 5px solid #2196F3 !important;
    }
    
    /* Spinner - Blue */
    .stSpinner > div {
        border-top-color: #2196F3 !important;
    }
    
    /* Horizontal rule */
    hr {
        border-color: #3d3d3d !important;
        margin: 2rem 0 !important;
    }
    
    /* Code blocks */
    code {
        background-color: #F5F5F5 !important;
        color: #1565C0 !important;
        padding: 0.2rem 0.4rem !important;
        border-radius: 4px !important;
    }
    
    /* Sidebar background */
    section[data-testid="stSidebar"] > div {
        background-color: #E3F2FD !important;
    }
    
    /* Column containers */
    [data-testid="column"] {
        background-color: transparent !important;
    }
    
    /* Links */
    a {
        color: #2196F3 !important;
    }
    a:hover {
        color: #1976D2 !important;
    }
    
    /* Chat messages */
    .stChatMessage {
        background-color: #F5F5F5 !important;
        border-radius: 8px !important;
        border: 1px solid #E0E0E0 !important;
        padding: 1rem !important;
        margin: 0.5rem 0 !important;
    }
    
    /* User chat message */
    [data-testid="stChatMessageContent"] {
        color: #212121 !important;
    }
    
    /* Chat input */
    [data-testid="stChatInput"] {
        background-color: #F5F5F5 !important;
        border: 2px solid #2196F3 !important;
        border-radius: 12px !important;
    }
    
    [data-testid="stChatInput"]:focus-within {
        border: 2px solid #1976D2 !important;
        box-shadow: 0 0 0 1px #1976D2 !important;
    }
    
    [data-testid="stChatInput"] > div {
        border-radius: 12px !important;
    }
    
    /* Fix white left border */
    [data-testid="stChatInput"] textarea {
        border-left: none !important;
    }
    
    /* Ensure chat message content is properly styled */
    [data-testid="stChatMessage"] > div {
        border-radius: 8px !important;
    }
    
    [data-testid="stChatInput"] textarea {
        background-color: #F5F5F5 !important;
        color: #212121 !important;
        border: none !important;
    }
    
    [data-testid="stChatInput"] textarea::placeholder {
        color: #757575 !important;
    }
    
    /* Chat input button */
    [data-testid="stChatInput"] button {
        background-color: #2196F3 !important;
        color: #FFFFFF !important;
    }
    
    /* Text area */
    textarea {
        background-color: #F5F5F5 !important;
        color: #212121 !important;
        border: none !important;
    }
    
    /* Select box styling */
    [data-testid="stSelectbox"] {
        background-color: #F5F5F5 !important;
        border-radius: 8px !important;
    }
    
    [data-testid="stSelectbox"] > div {
        border-radius: 8px !important;
    }
    
    [data-testid="stSelectbox"] > div > div {
        background-color: #F5F5F5 !important;
        color: #212121 !important;
        border: 1px solid #E0E0E0 !important;
        border-radius: 8px !important;
    }
    
    [data-testid="stSelectbox"] input {
        border-radius: 8px !important;
    }
    
    /* Dropdown arrow color */
    [data-testid="stSelectbox"] svg {
        fill: #2196F3 !important;
        color: #2196F3 !important;
    }
    
    /* Dropdown menu styling */
    [data-baseweb="popover"] {
        background-color: #FFFFFF !important;
        border-radius: 8px !important;
        border: 1px solid #E0E0E0 !important;
    }
    
    [data-baseweb="menu"] {
        background-color: #FFFFFF !important;
        border-radius: 8px !important;
    }
    
    [role="option"] {
        background-color: #FFFFFF !important;
        color: #212121 !important;
        border-radius: 4px !important;
    }
    
    [role="option"]:hover {
        background-color: #E3F2FD !important;
    }
    
    /* Radio button styling */
    [data-testid="stRadio"] label {
        color: #212121 !important;
    }
    
    [data-testid="stRadio"] > div {
        background-color: #F5F5F5 !important;
        padding: 0.5rem !important;
        border-radius: 6px !important;
    }
    
    /* Radio button input color - blue instead of red */
    [data-testid="stRadio"] input[type="radio"] {
        accent-color: #2196F3 !important;
        filter: hue-rotate(200deg) saturate(1.5) !important;
    }
    
    [data-testid="stRadio"] label span {
        accent-color: #2196F3 !important;
    }
    
    /* Force blue on radio container */
    [data-testid="stRadio"] [role="radiogroup"] {
        accent-color: #2196F3 !important;
    }
    
    /* Footer styling */
    footer {
        background-color: #F5F5F5 !important;
        color: #757575 !important;
        visibility: visible !important;
    }
    
    footer p {
        color: #757575 !important;
    }
    
    footer div {
        background-color: #F5F5F5 !important;
    }
    
    [data-testid="stBottomBlockContainer"] {
        background-color: #FFFFFF !important;
    }
    
    /* Bottom container */
    .main .block-container {
        padding-bottom: 5rem !important;
    }
    
    /* Streamlit footer */
    footer[data-testid="stFooter"] {
        background-color: #F5F5F5 !important;
    }
    
    footer[data-testid="stFooter"] > div {
        background-color: #F5F5F5 !important;
    }
    
    /* Chat footer area */
    [data-testid="stChatInputContainer"] {
        background-color: #FFFFFF !important;
        border-top: 1px solid #E0E0E0 !important;
    }
    
    /* More comprehensive table styling */
    div[data-testid="stDataFrame"] div[data-testid="StyledDataFrameContainer"] {
        background-color: transparent !important;
    }
    
    /* Remove grey box overlays */
    [data-testid="stDataFrame"] > div {
        background-color: transparent !important;
    }
    
    table {
        background-color: transparent !important;
        color: #212121 !important;
        border-collapse: collapse !important;
    }
    
    thead {
        background-color: #2196F3 !important;
    }
    
    th {
        background-color: #2196F3 !important;
        color: #FFFFFF !important;
    }
    
    td {
        background-color: #FAFAFA !important;
        color: #212121 !important;
    }
    
    tr:nth-child(even) td {
        background-color: #F5F5F5 !important;
    }
    
    /* Streamlit dataframe specific */
    .dataframe {
        background-color: transparent !important;
        color: #212121 !important;
    }
    
    .dataframe thead th {
        background-color: #2196F3 !important;
        color: #FFFFFF !important;
    }
    
    .dataframe tbody td {
        background-color: #FAFAFA !important;
        color: #212121 !important;
    }
    
    .dataframe tbody tr:nth-child(even) td {
        background-color: #F5F5F5 !important;
    }
    
    /* Override any remaining white backgrounds */
    div[class*="stDataFrame"] {
        background-color: transparent !important;
    }
</style>
""", unsafe_allow_html=True)

# Apply High Contrast Mode CSS if enabled
if st.session_state.high_contrast:
    st.markdown("""
    <style>
        /* High Contrast Mode Overrides */
        .main {
            background-color: #000000 !important;
        }
        
        .block-container {
            background-color: #000000 !important;
        }
        
        .stApp {
            background-color: #000000 !important;
        }
        
        header[data-testid="stHeader"] {
            background-color: #000000 !important;
        }
        
        [data-testid="stToolbar"] {
            background-color: #000000 !important;
        }
        
        .stApp > header {
            background-color: #000000 !important;
        }
        
        [data-testid="stVerticalBlock"] {
            background-color: #000000 !important;
        }
        
        /* Hide deploy button */
        [data-testid="stToolbar"] button[kind="header"] {
            display: none !important;
        }
        
        [data-testid="stSidebar"] {
            background-color: #000000 !important;
            border-right: 3px solid #FFFF00 !important;
        }
        
        section[data-testid="stSidebar"] > div {
            background-color: #000000 !important;
        }
        
        body, p, span, div, label, .stMarkdown, .stText {
            color: #FFFF00 !important;
            font-weight: 600 !important;
        }
        
        h1, h2, h3, h4, h5, h6 {
            color: #FFFF00 !important;
            text-shadow: 2px 2px 4px #000000 !important;
        }
        
        /* Main content headers must be yellow - all variants */
        .main h1, .main h2, .main h3, .main h4, .main h5, .main h6,
        .stMarkdown h1, .stMarkdown h2, .stMarkdown h3, .stMarkdown h4, .stMarkdown h5, .stMarkdown h6,
        [data-testid="stMarkdownContainer"] h1,
        [data-testid="stMarkdownContainer"] h2,
        [data-testid="stMarkdownContainer"] h3,
        [data-testid="stMarkdownContainer"] h4,
        [data-testid="stMarkdownContainer"] h5,
        [data-testid="stMarkdownContainer"] h6 {
            color: #FFFF00 !important;
            text-shadow: 2px 2px 4px #000000 !important;
            background-color: transparent !important;
        }
        
        /* Sidebar accessibility section - black text on yellow background */
        [data-testid="stSidebar"] h3:first-of-type {
            color: #000000 !important;
            background-color: #FFFF00 !important;
            padding: 0.5rem !important;
            border-radius: 4px !important;
            text-shadow: none !important;
        }
        
        button, .stButton button, .stDownloadButton button {
            background-color: #FFFF00 !important;
            color: #000000 !important;
            border: 3px solid #FFFFFF !important;
            font-weight: 700 !important;
        }
        
        button *, .stButton button *, .stDownloadButton button * {
            color: #000000 !important;
        }
        
        button:hover, .stButton button:hover {
            background-color: #FFFFFF !important;
            color: #000000 !important;
        }
        
        /* Radio buttons */
        [data-testid="stRadio"] {
            background-color: #000000 !important;
        }
        
        [data-testid="stRadio"] > div {
            background-color: #000000 !important;
            border: none !important;
        }
        
        [data-testid="stRadio"] label {
            color: #FFFF00 !important;
        }
        
        /* File uploader */
        [data-testid="stFileUploader"] {
            background-color: #000000 !important;
            border: 3px dashed #FFFF00 !important;
        }
        
        [data-testid="stFileUploader"] section {
            background-color: #000000 !important;
            border: 3px dashed #FFFF00 !important;
        }
        
        [data-testid="stFileUploader"] label {
            color: #FFFF00 !important;
        }
        
        [data-testid="stFileUploader"] small {
            color: #FFFF00 !important;
        }
        
        /* File uploader browse button */
        [data-testid="stFileUploader"] button {
            background-color: #2196F3 !important;
            color: #FFFFFF !important;
            border: 2px solid #1565C0 !important;
        }
        
        /* Selectbox */
        [data-testid="stSelectbox"] {
            background-color: #000000 !important;
        }
        
        [data-testid="stSelectbox"] > div > div {
            background-color: #000000 !important;
            color: #FFFF00 !important;
            border: 3px solid #FFFF00 !important;
        }
        
        [data-testid="stSelectbox"] label {
            color: #FFFF00 !important;
        }
        
        [data-testid="stSelectbox"] svg {
            fill: #FFFF00 !important;
            color: #FFFF00 !important;
        }
        
        /* Dropdown menu */
        [data-baseweb="popover"] {
            background-color: #000000 !important;
            border: 3px solid #FFFF00 !important;
        }
        
        [data-baseweb="menu"] {
            background-color: #000000 !important;
        }
        
        [role="option"] {
            background-color: #000000 !important;
            color: #FFFF00 !important;
        }
        
        [role="option"]:hover {
            background-color: #333333 !important;
            color: #FFFFFF !important;
        }
        
        /* Dataframes */
        [data-testid="stDataFrame"] {
            border: 3px solid #FFFF00 !important;
            background-color: #000000 !important;
        }
        
        [data-testid="stDataFrame"] thead tr th {
            background-color: #FFFF00 !important;
            color: #000000 !important;
        }
        
        [data-testid="stDataFrame"] tbody tr td {
            background-color: #000000 !important;
            color: #FFFF00 !important;
            border: 2px solid #FFFF00 !important;
        }
        
        table, .dataframe {
            background-color: #000000 !important;
            border: 3px solid #FFFF00 !important;
        }
        
        th {
            background-color: #FFFF00 !important;
            color: #000000 !important;
        }
        
        td {
            background-color: #000000 !important;
            color: #FFFF00 !important;
        }
        
        /* Input fields */
        input, textarea, select {
            background-color: #000000 !important;
            color: #FFFF00 !important;
            border: 3px solid #FFFF00 !important;
        }
        
        input::placeholder, textarea::placeholder {
            color: #CCCC00 !important;
        }
        
        /* Chat messages */
        .stChatMessage {
            background-color: #000000 !important;
            border: 3px solid #FFFF00 !important;
        }
        
        [data-testid="stChatMessageContent"] {
            color: #FFFF00 !important;
        }
        
        [data-testid="stChatInput"] {
            background-color: #000000 !important;
            border: 3px solid #FFFF00 !important;
        }
        
        [data-testid="stChatInput"] textarea {
            background-color: #000000 !important;
            color: #FFFF00 !important;
        }
        
        /* Metrics */
        [data-testid="stMetric"] {
            background-color: #000000 !important;
            border: 3px solid #FFFF00 !important;
        }
        
        [data-testid="stMetricValue"] {
            color: #FFFFFF !important;
        }
        
        [data-testid="stMetricLabel"] {
            color: #FFFF00 !important;
        }
        
        /* Expander */
        [data-testid="stExpander"] {
            background-color: #000000 !important;
            border: 3px solid #FFFF00 !important;
        }
        
        .streamlit-expanderHeader {
            background-color: #000000 !important;
            color: #FFFF00 !important;
        }
        
        /* Alerts */
        .stAlert, .stSuccess, .stInfo, .stWarning, .stError {
            background-color: #000000 !important;
            color: #FFFF00 !important;
            border: 3px solid #FFFF00 !important;
        }
        
        /* Links */
        a {
            color: #FFFF00 !important;
        }
        
        a:hover {
            color: #FFFFFF !important;
        }
        
        /* Footer */
        footer {
            background-color: #000000 !important;
            color: #FFFF00 !important;
        }
        
        footer p {
            color: #FFFF00 !important;
        }
        
        /* Sidebar text input (chat) in high contrast */
        [data-testid="stSidebar"] div[data-testid="stTextInput"] input {
            background-color: #000000 !important;
            border: 3px solid #FFFF00 !important;
            color: #FFFF00 !important;
        }
        
        [data-testid="stSidebar"] div[data-testid="stTextInput"] input:focus {
            border: 3px solid #FFFFFF !important;
            box-shadow: 0 0 0 1px #FFFFFF !important;
        }
        
        [data-testid="stSidebar"] div[data-testid="stTextInput"] input::placeholder {
            color: #CCCC00 !important;
        }
    </style>
    """, unsafe_allow_html=True)

# Apply Magnifier Mode CSS if enabled
if st.session_state.magnifier:
    st.markdown("""
    <style>
        /* Magnifier Mode - Larger Text */
        body, p, span, div, label, .stMarkdown, .stText {
            font-size: 1.3rem !important;
            line-height: 1.8 !important;
        }
        
        h1 {
            font-size: 3.5rem !important;
        }
        
        h2 {
            font-size: 2.8rem !important;
        }
        
        h3 {
            font-size: 2.2rem !important;
        }
        
        button {
            font-size: 1.4rem !important;
            padding: 1rem 2rem !important;
        }
        
        input, textarea {
            font-size: 1.3rem !important;
        }
        
        [data-testid="stMetricValue"] {
            font-size: 3rem !important;
        }
        
        [data-testid="stMetricLabel"] {
            font-size: 1.5rem !important;
        }
    </style>
    """, unsafe_allow_html=True)

# Text-to-Speech Function
def speak_text(text):
    """Generate text-to-speech audio"""
    if st.session_state.tts_enabled and text:
        try:
            from gtts import gTTS
            import base64
            import hashlib
            
            # Create unique key for this text to avoid duplicate TTS
            text_key = f"tts_{hashlib.md5(text[:100].encode()).hexdigest()}"
            
            # Only speak if we haven't already spoken this text
            if text_key not in st.session_state:
                st.session_state[text_key] = True
                
                # Generate speech
                tts = gTTS(text=text[:500], lang='en', slow=False)  # Limit to first 500 chars
                
                # Save to bytes
                audio_buffer = BytesIO()
                tts.write_to_fp(audio_buffer)
                audio_buffer.seek(0)
                
                # Convert to base64 for HTML audio player
                audio_base64 = base64.b64encode(audio_buffer.read()).decode()
                
                # Auto-play audio with controls visible
                st.markdown(f"""
                <audio controls autoplay style="width: 100%; margin: 10px 0;">
                    <source src="data:audio/mp3;base64,{audio_base64}" type="audio/mp3">
                </audio>
                """, unsafe_allow_html=True)
                return True
        except Exception as e:
            st.error(f"TTS Error: {str(e)}")
            return False
    return False

# Header with Accessibility Controls
st.markdown("""
<div style='margin-bottom: 1.5rem;'>
    <h1 style='margin: 0; color: #1565C0; font-size: 3.5rem; font-weight: 700;'>MedInsight</h1>
    <span style='color: #424242; font-size: 1.4rem; font-weight: 400;'>AI Tool to Help Personalise/Visualise RWE</span>
</div>
""", unsafe_allow_html=True)

st.markdown("**Accessibility**")
acc_col1, acc_col2, acc_col3 = st.columns(3)

with acc_col1:
    if st.button("Contrast" if not st.session_state.high_contrast else "✓ Contrast", key="hc_btn", use_container_width=True):
        st.session_state.high_contrast = not st.session_state.high_contrast
        st.rerun()

with acc_col2:
    if st.button("Magnify" if not st.session_state.magnifier else "✓ Magnify", key="mag_btn", use_container_width=True):
        st.session_state.magnifier = not st.session_state.magnifier
        st.rerun()

with acc_col3:
    if st.button("TTS" if not st.session_state.tts_enabled else "TTS ON", key="tts_btn", use_container_width=True):
        st.session_state.tts_enabled = not st.session_state.tts_enabled
        st.rerun()

st.markdown("---")

# Helper functions
def extract_text_from_pptx(file):
    """Extract text and tables from PowerPoint"""
    if not PPTX_SUPPORT:
        return "PowerPoint support not available. Install python-pptx.", []
    
    prs = Presentation(file)
    text_content = []
    tables_data = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
            
            if shape.has_table:
                table = shape.table
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                if table_data:
                    tables_data.append({"slide": slide_num, "data": table_data})
        
        if slide_text:
            text_content.append(f"Slide {slide_num}:\n" + "\n".join(slide_text))
    
    return "\n\n".join(text_content), tables_data

def parse_json_to_dataframe(file):
    """Parse JSON file to DataFrame"""
    try:
        content = json.load(file)
        if isinstance(content, list):
            return pd.DataFrame(content)
        elif isinstance(content, dict):
            for key, value in content.items():
                if isinstance(value, list) and len(value) > 0:
                    return pd.DataFrame(value)
            return pd.DataFrame([content])
        return pd.DataFrame()
    except Exception as e:
        st.error(f"Error parsing JSON: {str(e)}")
        return pd.DataFrame()

# Data Input Section
st.markdown("#### Choose input method:")
col1, col2 = st.columns(2)

with col1:
    if st.button("Upload File", use_container_width=True, type="primary" if st.session_state.input_method == "Upload File" else "secondary"):
        st.session_state.input_method = "Upload File"
        st.rerun()

with col2:
    if st.button("Enter Text Directly", use_container_width=True, type="primary" if st.session_state.input_method == "Enter Text Directly" else "secondary"):
        st.session_state.input_method = "Enter Text Directly"
        st.rerun()

input_method = st.session_state.input_method

uploaded_file = None
text_input = None
file_type = None

if input_method == "Upload File":
    file_type = st.selectbox(
        "Select file type:",
        ["Data File (Excel/CSV)", "PDF Document", "PowerPoint Presentation", "Text Document", "JSON Data"]
    )
    
    if file_type == "Data File (Excel/CSV)":
        uploaded_file = st.file_uploader(
            "Choose an Excel or CSV file",
            type=['xlsx', 'xls', 'csv'],
            help="Upload your data file for analysis"
        )
    elif file_type == "PDF Document":
        uploaded_file = st.file_uploader(
            "Choose a PDF file",
            type=['pdf'],
            help="Upload PDF document for analysis"
        )
    elif file_type == "PowerPoint Presentation":
        uploaded_file = st.file_uploader(
            "Choose a PowerPoint file",
            type=['pptx'],
            help="Upload PowerPoint with data tables or analysis"
        )
    elif file_type == "Text Document":
        uploaded_file = st.file_uploader(
            "Choose a text file",
            type=['txt', 'md'],
            help="Upload text document with analysis or data"
        )
    else:  # JSON Data
        uploaded_file = st.file_uploader(
            "Choose a JSON file",
            type=['json'],
            help="Upload JSON file with structured data"
        )
else:
    # Mobile-friendly text input with submit button
    st.markdown("#### Enter Text (Mobile Supported)")
    text_input = st.text_area(
        "Enter your data or analysis:",
        height=200,
        placeholder="Paste your text, data, or analysis here...\n\nYou can paste:\n- CSV data\n- Analysis text\n- Research findings\n- Any text for AI analysis",
        help="Enter any text for AI analysis and insights",
        key="text_input_area"
    )
    
    # Add submit button for mobile users
    col1, col2 = st.columns([3, 1])
    with col1:
        submit_text = st.button("Submit Text", type="primary", use_container_width=True, help="Click to process your text")
    with col2:
        clear_text = st.button("Clear", use_container_width=True, help="Clear text input")
    
    if clear_text:
        st.session_state.text_input_area = ""
        st.rerun()
    
    # Only process if submit button is clicked or text exists
    if not submit_text and not text_input:
        text_input = None
    
    file_type = "Text Input"

if uploaded_file is not None or (text_input and text_input.strip()):
    # User Profile Section
    st.markdown("---")
    st.markdown("### 👤 User Profile")
    st.markdown("Help us customize the analysis for you:")
    
    col1, col2, col3 = st.columns(3)
    
    with col1:
        knowledge_level = st.selectbox(
            "Knowledge Level",
            [
                "Beginner - New to RWE/Clinical Data",
                "Intermediate - Some experience with data analysis",
                "Advanced - Experienced with RWE studies",
                "Expert - Clinical research professional"
            ],
            help="Select your familiarity with Real-World Evidence and clinical data"
        )
    
    with col2:
        country = st.selectbox(
            "Country/Language",
            [
                "United States - English",
                "United Kingdom - English",
                "Canada - English/French",
                "Australia - English",
                "Germany - German",
                "France - French",
                "Spain - Spanish",
                "Italy - Italian",
                "Japan - Japanese",
                "China - Mandarin",
                "Brazil - Portuguese",
                "India - English/Hindi",
                "Other"
            ],
            help="Select your country and preferred language"
        )
    
    with col3:
        hcp_type = st.selectbox(
            "Healthcare Professional Type",
            [
                "Physician - General Practice",
                "Physician - Specialist",
                "Nurse Practitioner",
                "Pharmacist",
                "Researcher/Scientist",
                "Clinical Trial Coordinator",
                "Healthcare Administrator",
                "Medical Affairs",
                "Regulatory Affairs",
                "Data Analyst/Statistician",
                "Student/Trainee",
                "Other"
            ],
            help="Select your role in healthcare"
        )
    
    # Store in session state for AI to use
    st.session_state.user_profile = {
        "knowledge_level": knowledge_level.split(" - ")[0],
        "knowledge_desc": knowledge_level,
        "country": country.split(" - ")[0],
        "language": country.split(" - ")[-1] if " - " in country else "English",
        "hcp_type": hcp_type.split(" - ")[0] if " - " in hcp_type else hcp_type
    }
    
    st.markdown("---")
    
    try:
        df = None
        extracted_text = None
        
        # Process based on input type
        if text_input and text_input.strip():
            extracted_text = text_input.strip()
            st.success(f"Text input received: {len(extracted_text)} characters")
            
            # Try to parse as CSV
            try:
                from io import StringIO
                df = pd.read_csv(StringIO(extracted_text))
                st.info("Detected CSV format in text - parsed as DataFrame")
            except:
                pass
        
        elif uploaded_file is not None:
            if file_type == "Data File (Excel/CSV)":
                with st.spinner("Loading data..."):
                    if uploaded_file.name.endswith('.csv'):
                        try:
                            df = pd.read_csv(uploaded_file, encoding='utf-8')
                        except UnicodeDecodeError:
                            uploaded_file.seek(0)
                            df = pd.read_csv(uploaded_file, encoding='latin-1')
                    else:
                        df = pd.read_excel(uploaded_file)
                st.success(f"Loaded data: {len(df)} rows x {len(df.columns)} columns")
            
            elif file_type == "PDF Document":
                if PDF_SUPPORT:
                    with st.spinner("Extracting text from PDF..."):
                        try:
                            pdf_reader = PdfReader(uploaded_file)
                            extracted_text = ""
                            for page_num, page in enumerate(pdf_reader.pages, 1):
                                extracted_text += f"\n--- Page {page_num} ---\n"
                                extracted_text += page.extract_text()
                            
                            st.success(f"Extracted text from {len(pdf_reader.pages)} pages")
                            st.markdown("### PDF Content")
                            with st.expander("View extracted text", expanded=False):
                                st.text_area("PDF Content", extracted_text, height=300, key="pdf_text")
                        except Exception as e:
                            st.error(f"Error extracting PDF: {str(e)}")
                else:
                    st.error("PDF support not available. Install PyPDF2 package.")
            
            elif file_type == "PowerPoint Presentation":
                extracted_text, tables = extract_text_from_pptx(uploaded_file)
                st.markdown("### Extracted Content")
                with st.expander("View extracted text", expanded=False):
                    st.text_area("PowerPoint Content", extracted_text, height=300, key="pptx_text")
                
                if tables:
                    st.markdown("### Extracted Tables")
                    for idx, table_info in enumerate(tables):
                        st.markdown(f"**Table from Slide {table_info['slide']}**")
                        table_data = table_info['data']
                        if len(table_data) > 1:
                            df_temp = pd.DataFrame(table_data[1:], columns=table_data[0])
                            st.dataframe(df_temp, use_container_width=True)
                            if df is None:
                                df = df_temp
            
            elif file_type == "Text Document":
                try:
                    extracted_text = uploaded_file.read().decode('utf-8')
                except UnicodeDecodeError:
                    uploaded_file.seek(0)
                    extracted_text = uploaded_file.read().decode('latin-1')
                st.markdown("### Document Content")
                with st.expander("View document text", expanded=False):
                    st.text_area("Text Content", extracted_text, height=300, key="text_content")
            
            elif file_type == "JSON Data":
                df = parse_json_to_dataframe(uploaded_file)
                if not df.empty:
                    st.success(f"Loaded JSON data: {len(df)} records")
        
        # Display data preview
        st.markdown("### Data Preview")
        st.markdown("**Quick overview of your data:** The metrics below show you the size and structure of your dataset. Click the expandable sections to explore the details.")
        
        if df is not None and not df.empty:
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("Total Rows", f"{len(df):,}")
            with col2:
                st.metric("Total Columns", len(df.columns))
            with col3:
                numeric_cols = df.select_dtypes(include=['number']).columns
                st.metric("Numeric Columns", len(numeric_cols))
            
            with st.expander("View Data Sample", expanded=True):
                st.dataframe(df.head(10), use_container_width=True)
            
            # Column information
            with st.expander("Column Details"):
                col_info = pd.DataFrame({
                    'Column Name': df.columns,
                    'Data Type': [str(dtype) for dtype in df.dtypes.values],
                    'Non-Null Count': df.count().values,
                    'Null Count': df.isnull().sum().values,
                    'Unique Values': [df[col].nunique() for col in df.columns]
                })
                st.dataframe(col_info, use_container_width=True)
        
        elif extracted_text:
            col1, col2 = st.columns(2)
            with col1:
                st.metric("Characters", f"{len(extracted_text):,}")
            with col2:
                st.metric("Words", f"{len(extracted_text.split()):,}")
            
            with st.expander("View Content Preview", expanded=True):
                st.markdown(extracted_text[:1000] + ("..." if len(extracted_text) > 1000 else ""))
        
        st.markdown("---")
        
        # ============================================
        # AI CHAT ASSISTANT - AVAILABLE AFTER DATA UPLOAD
        # ============================================
        st.markdown("### AI Chat Assistant")
        st.markdown("**Ask questions about your data** - Get insights, clarification, and analysis help. The AI has context about your uploaded data.")

        # Example question buttons
        st.markdown("**Quick Start:**")
        example_cols = st.columns(4)
        
        example_questions = [
            "What are the key trends in this data?",
            "Summarize the main findings",
            "What patterns should I investigate?",
            "Explain the data structure"
        ]
        
        if 'example_question_clicked' not in st.session_state:
            st.session_state.example_question_clicked = None
        
        for idx, col in enumerate(example_cols):
            with col:
                if st.button(example_questions[idx], key=f"example_q_{idx}", use_container_width=True):
                    st.session_state.example_question_clicked = example_questions[idx]
                    st.rerun()
        
        st.markdown("")

        # Chat interface
        chat_col1, chat_col2 = st.columns([4, 1])

        with chat_col1:
            general_prompt = st.chat_input("Ask AI anything about your data or analysis...")

        with chat_col2:
            if st.session_state.get('chatbot_messages', []):
                if st.button("Clear Chat", type="secondary", use_container_width=True, key="clear_general_chat"):
                    st.session_state.chatbot_messages = []
                    st.rerun()

        # Display chat messages
        if st.session_state.get('chatbot_messages', []):
            with st.container():
                for message in st.session_state.chatbot_messages:
                    with st.chat_message(message["role"]):
                        st.markdown(message["content"], unsafe_allow_html=True)

        # Handle new messages from general chat (including example questions)
        active_prompt = general_prompt if general_prompt and general_prompt.strip() else st.session_state.example_question_clicked
        
        if active_prompt:
            # Check if this message was already processed
            if 'last_general_prompt' not in st.session_state or st.session_state.last_general_prompt != active_prompt:
                st.session_state.last_general_prompt = active_prompt
                st.session_state.example_question_clicked = None  # Reset example question
                
                # Add user message to chat history
                st.session_state.chatbot_messages.append({"role": "user", "content": active_prompt})
                
                # Build context from analysis if available
                context = ""
                if 'analysis_text' in st.session_state and st.session_state.analysis_text:
                    context += f"\n\nPrevious Analysis:\n{st.session_state.analysis_text[:1000]}"
                
                if 'df' in st.session_state and st.session_state.get('df') is not None:
                    df_context = st.session_state['df']
                    context += f"\n\nData Context:\n"
                    context += f"- Total Rows: {len(df_context)}\n"
                    context += f"- Total Columns: {len(df_context.columns)}\n"
                    context += f"- Columns: {', '.join(df_context.columns.tolist()[:15])}\n\n"
                    
                    # Add statistical summary for numeric columns
                    numeric_cols = df_context.select_dtypes(include=['number']).columns[:5]
                    if len(numeric_cols) > 0:
                        context += "Key Numeric Statistics:\n"
                        for col in numeric_cols:
                            context += f"  - {col}: mean={df_context[col].mean():.2f}, median={df_context[col].median():.2f}, std={df_context[col].std():.2f}\n"
                    
                    # Add value counts for categorical columns
                    categorical_cols = df_context.select_dtypes(include=['object']).columns[:3]
                    if len(categorical_cols) > 0:
                        context += "\nKey Categorical Distributions:\n"
                        for col in categorical_cols:
                            top_values = df_context[col].value_counts().head(5).to_dict()
                            context += f"  - {col}: {top_values}\n"
                    
                    # Add sample rows
                    context += f"\nSample Data (first 3 rows):\n{df_context.head(3).to_string()}"
                
                # Create HCP-focused prompt with user profile
                user_profile = st.session_state.get('user_profile', {})
                knowledge = user_profile.get('knowledge_level', 'Intermediate')
                language = user_profile.get('language', 'English')
                hcp_role = user_profile.get('hcp_type', 'Healthcare Professional')
                
                full_prompt = f"""You are MedInsight, an AI assistant helping Healthcare Professionals (HCPs) analyze Real-World Evidence (RWE) data.

User Profile:
- Role: {hcp_role}
- Knowledge Level: {knowledge}
- Language: {language}

User Question: {active_prompt}

DATA AVAILABLE TO YOU:
{context}

INSTRUCTIONS:
Answer the user's question based SPECIFICALLY on the data provided above. Use actual numbers, column names, and values from the dataset.
- Reference specific findings from the data (e.g., "The average weight change was -11.3 kg in the Mounjaro group")
- Cite actual statistics and distributions shown in the data
- Focus on the treatments, outcomes, and patient characteristics present in THIS dataset
- Tailor complexity to {knowledge} level
- If the data doesn't contain information to answer the question, acknowledge this and suggest related insights from available data

Keep your response concise (3-4 paragraphs max), clinically relevant, and evidence-based using the actual data provided.

Response:"""
                
                # Generate AI response
                try:
                    with st.spinner("AI is thinking..."):
                        response = generate_ai_response(full_prompt)
                    
                    # Add assistant response to chat history
                    st.session_state.chatbot_messages.append({"role": "assistant", "content": response})
                    
                    # Auto-play TTS for chatbot response if enabled
                    if st.session_state.tts_enabled:
                        speak_text(response)
                    
                    # Rerun to display new messages
                    st.rerun()
                        
                except Exception as e:
                    st.error(f"Error generating response: {str(e)}")
        
        st.markdown("---")
        
        # Generate summary with AI
        st.markdown("### AI Analysis")
        
        # Add explanation for non-technical users
        with st.expander("What does AI Analysis do?", expanded=False):
            st.markdown("""
            **AI Analysis helps you understand your data without needing technical expertise.**
            
            **What it does:**
            - Reads through your data or document automatically
            - Identifies important patterns and trends
            - Explains findings in plain language
            - Provides practical recommendations
            
            **Why it's valuable:**
            - Saves hours of manual data review
            - Highlights insights you might miss
            - Tailored to your expertise level
            - Focuses on clinically relevant findings
            
            **How to use it:**
            Simply click the button below and the AI will analyze your data and present the findings in an easy-to-understand format.
            """)
        
        if st.button("Ask AI to Analyze", type="primary", use_container_width=True):
            if not model:
                st.error("⚠️ API key not configured. Please set GEMINI_API_KEY in .env file")
            else:
                with st.spinner("AI is analyzing your data..."):
                    # Prepare data summary - prioritize extracted text for PowerPoint/Text files
                    if extracted_text and file_type in ["PowerPoint Presentation", "Text Document"]:
                        # For PowerPoint and text files, analyze the content, not the tables
                        user_profile = st.session_state.get('user_profile', {})
                        knowledge = user_profile.get('knowledge_level', 'Intermediate')
                        language = user_profile.get('language', 'English')
                        hcp_role = user_profile.get('hcp_type', 'Healthcare Professional')
                        
                        # Limit text content to prevent token overflow (max 3000 characters)
                        max_chars = 3000
                        truncated_text = extracted_text[:max_chars]
                        if len(extracted_text) > max_chars:
                            truncated_text += "...\n[Content truncated for analysis]"
                        
                        data_summary = f"""
Content Analysis:
- Type: {file_type}
- Length: {len(extracted_text)} characters
- Analyzing: First {min(len(extracted_text), max_chars)} characters

Content:
{truncated_text}
"""
                        prompt = f"""You are analyzing this content for a {hcp_role} with {knowledge} knowledge level. 
Language preference: {language} (adapt complexity and terminology accordingly).

Personalize your analysis for their clinical perspective and provide:

1. Executive Summary: Key takeaways relevant to clinical practice
2. Clinical Insights: How this information applies to patient care and treatment decisions
3. Actionable Recommendations: Specific steps HCPs can take based on this information
4. Key Concepts Explained: Medical/pharmaceutical concepts in practical terms
5. Evidence & Data Quality: Assessment of the strength and relevance of the information
6. Practice Implications: How this impacts day-to-day healthcare delivery

Present the analysis in a clear, professional format that an HCP would find valuable for clinical decision-making.

{data_summary}
"""
                    elif df is not None and not df.empty:
                        user_profile = st.session_state.get('user_profile', {})
                        knowledge = user_profile.get('knowledge_level', 'Intermediate')
                        language = user_profile.get('language', 'English')
                        hcp_role = user_profile.get('hcp_type', 'Healthcare Professional')
                        
                        # Limit data to prevent token overflow
                        sample_size = min(5, len(df))
                        cols_to_show = min(10, len(df.columns))
                        
                        data_summary = f"""
Dataset Overview:
- Total Rows: {len(df)}
- Total Columns: {len(df.columns)}
- Column Names: {', '.join(df.columns.tolist()[:cols_to_show])}{'...' if len(df.columns) > cols_to_show else ''}

Data Types (first {cols_to_show} columns):
{df.dtypes.head(cols_to_show).to_string()}

Statistical Summary (first {cols_to_show} numeric columns):
{df.describe().iloc[:, :cols_to_show].to_string()}

Sample Data (first {sample_size} rows, first {cols_to_show} columns):
{df.head(sample_size).iloc[:, :cols_to_show].to_string()}
"""
                        prompt = f"""You are analyzing this Real-World Evidence (RWE) data for a {hcp_role} with {knowledge} knowledge level.
Language preference: {language} (adapt complexity and terminology accordingly).

Personalize your analysis for clinical practice:

1. Clinical Summary: What does this data tell us about patient outcomes, treatment effectiveness, or healthcare patterns?
2. Patient Impact: How do these findings relate to patient care and treatment decisions?
3. Evidence Quality: Assess the data quality, sample size, and statistical significance
4. Key Metrics: Highlight the most important clinical measures and their implications
5. Actionable Insights: Specific recommendations for HCPs based on this RWE
6. Comparative Analysis: How do these findings compare to established clinical standards or previous studies?
7. Practice Implementation: Practical steps HCPs can take to apply these insights

Present insights in a format that supports evidence-based clinical decision-making.

{data_summary}
"""
                    else:
                        data_summary = "No data available."
                        prompt = f"""No data available to analyze.

{data_summary}
"""
                    
                    try:
                        analysis_text = generate_ai_response(prompt)
                        st.session_state.analysis_text = analysis_text
                        
                        # Display results with clear user-friendly structure
                        st.markdown("---")
                        st.markdown("### Your Analysis Results")
                        
                        st.info("**Understanding Your Results**: The AI has reviewed your data and broken down the findings into clear sections below. Each section focuses on a specific aspect to help you make informed decisions.")
                        
                        # Display the analysis in a well-structured format
                        st.markdown(analysis_text)
                        
                        # Add helpful context after analysis
                        st.markdown("---")
                        with st.expander("How to Use These Results", expanded=False):
                            st.markdown("""
                            **Making the Most of Your Analysis:**
                            
                            1. **Read the Clinical Summary first** - This gives you the big picture
                            2. **Review Patient Impact** - See how findings affect care decisions
                            3. **Check Actionable Insights** - Get specific recommendations you can implement
                            4. **Use the buttons below** to:
                               - Generate a professional summary for reports or presentations
                               - Create visual charts and graphs to see patterns
                               - Ask follow-up questions in the chatbot
                            
                            **Next Steps:**
                            - Save or download the summary for your records
                            - Share findings with your team
                            - Use visualizations to present data to stakeholders
                            """)
                        
                        # Auto-speak analysis if TTS is enabled
                        if st.session_state.tts_enabled:
                            speak_text(analysis_text)
                        
                        # Store analysis in session state for visualization
                        st.session_state['analysis_done'] = True
                        st.session_state['analysis_text'] = analysis_text
                        st.session_state['df'] = df
                        
                        # ============================================
                        # CHAT INTERFACE - AFTER ANALYSIS
                        # ============================================
                        st.markdown("---")
                        st.markdown("### Ask Follow-Up Questions")
                        st.markdown("**Have questions about your analysis?** Chat with AI to dive deeper into specific findings or get clarification.")
                        
                        # Example question buttons for analysis follow-up
                        st.markdown("**Quick Questions:**")
                        followup_cols = st.columns(4)
                        
                        followup_questions = [
                            "Explain the clinical significance",
                            "What are the limitations?",
                            "Compare to typical outcomes",
                            "What should I investigate further?"
                        ]
                        
                        if 'followup_question_clicked' not in st.session_state:
                            st.session_state.followup_question_clicked = None
                        
                        for idx, col in enumerate(followup_cols):
                            with col:
                                if st.button(followup_questions[idx], key=f"followup_q_{idx}", use_container_width=True):
                                    st.session_state.followup_question_clicked = followup_questions[idx]
                                    st.rerun()
                        
                        st.markdown("")
                        
                        with st.expander("More ways to use the chat", expanded=False):
                            st.markdown("""
                            **Ask the AI follow-up questions about your analysis:**
                            
                            - Get clarification on specific findings
                            - Request alternative interpretations
                            - Ask for more detail on any section
                            - Compare findings to other studies or standards
                            
                            **Example questions:**
                            - "Can you explain the clinical significance of [finding] in more detail?"
                            - "What are the limitations of this analysis?"
                            - "How does this compare to typical outcomes?"
                            - "What should I investigate further?"
                            """)
                        
                        # Chat interface
                        chat_col1, chat_col2 = st.columns([4, 1])
                        
                        with chat_col1:
                            prompt = st.chat_input("Ask AI about your analysis...")
                        
                        with chat_col2:
                            if st.session_state.get('chatbot_messages', []):
                                if st.button("Clear Chat", type="secondary", use_container_width=True, key="clear_analysis_chat"):
                                    st.session_state.chatbot_messages = []
                                    st.rerun()
                        
                        # Display chat messages
                        if st.session_state.get('chatbot_messages', []):
                            chat_container = st.container()
                            with chat_container:
                                for message in st.session_state.chatbot_messages:
                                    with st.chat_message(message["role"]):
                                        st.markdown(message["content"], unsafe_allow_html=True)
                        
                        # Handle new messages (including example questions)
                        active_followup = prompt if prompt and prompt.strip() else st.session_state.followup_question_clicked
                        
                        if active_followup:
                            # Check if this message was already processed
                            if 'last_prompt' not in st.session_state or st.session_state.last_prompt != active_followup:
                                st.session_state.last_prompt = active_followup
                                st.session_state.followup_question_clicked = None  # Reset example question
                                
                                # Add user message to chat history
                                st.session_state.chatbot_messages.append({"role": "user", "content": active_followup})
                                
                                # Build context from analysis
                                context = f"\n\nPrevious Analysis:\n{analysis_text[:1000]}"
                                
                                if df is not None:
                                    context += f"\n\nData Context: {len(df)} rows, {len(df.columns)} columns"
                                    context += f"\nColumns: {', '.join(df.columns.tolist()[:10])}"
                                
                                # Create HCP-focused prompt with user profile
                                user_profile = st.session_state.get('user_profile', {})
                                knowledge = user_profile.get('knowledge_level', 'Intermediate')
                                language = user_profile.get('language', 'English')
                                hcp_role = user_profile.get('hcp_type', 'Healthcare Professional')
                                
                                full_prompt = f"""You are MedInsight, an AI assistant helping Healthcare Professionals (HCPs) analyze Real-World Evidence (RWE) data.

User Profile:
- Role: {hcp_role}
- Knowledge Level: {knowledge}
- Language: {language}

User Question: {active_followup}

{context}

Provide a helpful, professional response tailored for this {hcp_role} with {knowledge} knowledge level. 
Adapt your response complexity and terminology to match their expertise. Focus on:
- Clinical relevance and practical application for their specific role
- Evidence-based insights at the appropriate technical level
- Clear, actionable guidance
- Terminology suitable for {knowledge} level

Response:"""
                                
                                # Generate AI response
                                try:
                                    with st.spinner("AI is thinking..."):
                                        response = generate_ai_response(full_prompt)
                                    
                                    # Add assistant response to chat history
                                    st.session_state.chatbot_messages.append({"role": "assistant", "content": response})
                                    
                                    # Auto-play TTS for chatbot response if enabled
                                    if st.session_state.tts_enabled:
                                        speak_text(response)
                                    
                                    # Rerun to display new messages
                                    st.rerun()
                                        
                                except Exception as e:
                                    st.error(f"Error generating response: {str(e)}")
                        
                    except Exception as e:
                        st.error(f"Error generating analysis: {str(e)}")
        
        # Show visualization and professional summary buttons after analysis is done
        if st.session_state.get('analysis_done', False):
            st.markdown("---")
            st.markdown("### Next Actions")
            st.markdown("**Choose what you'd like to do with your analysis:**")
            
            # Professional Summary Generator
            col1, col2 = st.columns(2)
            
            with col1:
                st.markdown("**Professional Summary**")
                st.caption("Creates a formal report suitable for presentations, publications, or sharing with colleagues. Includes executive summary and structured findings.")
                if st.button("Generate Professional Summary", type="primary", use_container_width=True):
                    with st.spinner("Creating professional summary..."):
                        # Build comprehensive context
                        summary_context = ""
                        
                        if 'analysis_text' in st.session_state and st.session_state.analysis_text:
                            summary_context += f"Analysis:\n{st.session_state.analysis_text}\n\n"
                        
                        if df is not None and not df.empty:
                            summary_context += f"Dataset: {len(df)} rows, {len(df.columns)} columns\n"
                            summary_context += f"Columns: {', '.join(df.columns.tolist())}\n"
                        
                        # Generate professional summary
                        summary_prompt = f"""Create a concise, professional executive summary of this Real-World Evidence (RWE) analysis for Healthcare Professionals.

Context:
{summary_context}

Format your summary with:

**Executive Summary**
[2-3 sentence overview of key findings]

**Key Insights**
• [Most important finding 1]
• [Most important finding 2]
• [Most important finding 3]

**Clinical Implications**
[How these findings impact patient care and treatment decisions]

**Recommendations**
1. [Actionable recommendation 1]
2. [Actionable recommendation 2]
3. [Actionable recommendation 3]

**Data Quality & Limitations**
[Brief assessment of data reliability and any caveats]

Keep it professional, concise, and focused on actionable insights for HCPs."""
                        
                        try:
                            professional_summary = generate_ai_response(summary_prompt)
                            st.session_state['professional_summary'] = professional_summary
                            
                            st.markdown("### 📋 Professional Summary")
                            st.markdown(professional_summary)
                            
                            # Auto-speak summary if TTS is enabled
                            if st.session_state.tts_enabled:
                                speak_text(professional_summary)
                            
                            # Download option for summary
                            summary_report = f"""Professional RWE Analysis Summary
{'=' * 50}

Generated: {pd.Timestamp.now().strftime('%Y-%m-%d %H:%M:%S')}
File: {uploaded_file.name if uploaded_file else 'Text Input'}

{professional_summary}

{'=' * 50}
Powered by MedInsight AI
"""
                            st.download_button(
                                label="⬇️ Download Summary (TXT)",
                                data=summary_report.encode('utf-8'),
                                file_name="professional_summary.txt",
                                mime="text/plain",
                                use_container_width=True
                            )
                        except Exception as e:
                            st.error(f"Error generating summary: {str(e)}")
            
            with col2:
                # Show existing professional summary if available
                if 'professional_summary' in st.session_state:
                    if st.button("📋 View Saved Summary", use_container_width=True):
                        st.markdown("### 📋 Professional Summary")
                        st.markdown(st.session_state['professional_summary'])
            
            st.markdown("---")
            
            # Only show visualisation button if we have a dataframe
            if df is not None and not df.empty:
                st.markdown("**Visual Insights**")
                st.caption("Transforms your data into interactive charts and 3D graphs. Visual representations make it easier to spot trends, compare values, and present findings to others.")
                visualize_button = st.button("Generate Visualizations & 3D AR Graphs", use_container_width=True)
                if visualize_button:
                    # Show AI Analysis again
                    st.markdown("### AI Analysis")
                    st.markdown(st.session_state.get('analysis_text', ''))
                    st.markdown("---")
                    with st.spinner("Creating interactive visualisations..."):
                        st.markdown("### Data Visualizations")
                        
                        st.info("**Understanding Your Charts**: Each visualization below shows a different aspect of your data. You can interact with the charts - hover over points for details, zoom in/out, and rotate 3D graphs. These visuals help you see patterns that might not be obvious in raw numbers.")
                        
                        with st.expander("Chart Guide - What Each Visualization Shows", expanded=False):
                            st.markdown("""
                            **3D Scatter Plot**: Shows relationships between three variables at once. Each dot represents a data point, and colors group similar items together.
                            
                            **Distribution Charts**: Show how values are spread out - helping you see if data is clustered or evenly distributed.
                            
                            **Bar Charts**: Compare categories side-by-side to quickly see which groups have higher or lower values.
                            
                            **Correlation Matrix**: Shows how variables relate to each other. Blue colors indicate positive relationships, while lighter colors show little to no relationship.
                            
                            **How to use them:**
                            - Hover your mouse over any element for detailed information
                            - Click and drag to rotate 3D charts
                            - Use the toolbar icons to zoom, pan, or download images
                            - Take screenshots for reports or presentations
                            """)
                        
                        # Get numeric and categorical columns
                        numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
                        categorical_cols = df.select_dtypes(include=['object', 'category']).columns.tolist()
                    
                    if len(numeric_cols) >= 2:
                        # 3D Scatter Plot (AR-ready)
                        st.markdown("---")
                        st.markdown("#### 3D Interactive Scatter Plot")
                        st.markdown("**What this shows:** A three-dimensional view of your data where each point represents a record. This helps you see clusters, outliers, and relationships between multiple variables at once. Try rotating it with your mouse!")
                        
                        import plotly.graph_objects as go
                        
                        # Select up to 3 numeric columns for 3D plot
                        x_col = numeric_cols[0]
                        y_col = numeric_cols[1] if len(numeric_cols) > 1 else numeric_cols[0]
                        z_col = numeric_cols[2] if len(numeric_cols) > 2 else numeric_cols[0]
                        
                        # Color by categorical if available
                        color_col = categorical_cols[0] if categorical_cols else None
                        
                        if color_col:
                            # Create color mapping for categorical column
                            unique_categories = df[color_col].unique()
                            colors = ['#2196F3' if i % 2 == 0 else '#1565C0' for i in range(len(unique_categories))]
                            color_map = {cat: colors[i] for i, cat in enumerate(unique_categories)}
                            
                            fig = go.Figure(data=[go.Scatter3d(
                                x=df[x_col],
                                y=df[y_col],
                                z=df[z_col],
                                mode='markers',
                                marker=dict(
                                    size=6,
                                    color=[color_map[val] for val in df[color_col]],
                                    line=dict(color='#1976D2', width=0.5)
                                ),
                                text=[f"{color_col}: {val}" for val in df[color_col]],
                                hovertemplate=f'<b>%{{text}}</b><br>{x_col}: %{{x}}<br>{y_col}: %{{y}}<br>{z_col}: %{{z}}<extra></extra>'
                            )])
                        else:
                            fig = go.Figure(data=[go.Scatter3d(
                                x=df[x_col],
                                y=df[y_col],
                                z=df[z_col],
                                mode='markers',
                                marker=dict(
                                    size=6,
                                    color='#2196F3',
                                    line=dict(color='#1976D2', width=0.5)
                                ),
                                hovertemplate=f'{x_col}: %{{x}}<br>{y_col}: %{{y}}<br>{z_col}: %{{z}}<extra></extra>'
                            )])
                        
                        fig.update_layout(
                            title=dict(
                                text=f'3D Scatter: {x_col} vs {y_col} vs {z_col}',
                                font=dict(color='#1565C0', size=18)
                            ),
                            scene=dict(
                                xaxis=dict(title=x_col, gridcolor='#E3F2FD', backgroundcolor='#FFFFFF'),
                                yaxis=dict(title=y_col, gridcolor='#E3F2FD', backgroundcolor='#FFFFFF'),
                                zaxis=dict(title=z_col, gridcolor='#E3F2FD', backgroundcolor='#FFFFFF'),
                            ),
                            paper_bgcolor='#FFFFFF',
                            font=dict(color='#212121'),
                            height=600
                        )
                        
                        st.plotly_chart(fig, use_container_width=True)
                    
                    # Distribution plots
                    if numeric_cols:
                        st.markdown("---")
                        st.markdown("#### Distribution Analysis")
                        st.markdown("**What this shows:** How your numeric values are spread out. Tall bars indicate common values, while short bars show rare values. This helps identify typical ranges and unusual data points.")
                        
                        for col in numeric_cols[:3]:  # Show first 3 numeric columns
                            fig = go.Figure(data=[go.Histogram(
                                x=df[col],
                                marker_color='#2196F3',
                                marker_line_color='#1565C0',
                                marker_line_width=1.5,
                                opacity=0.85
                            )])
                            
                            fig.update_layout(
                                title=dict(
                                    text=f'Distribution of {col}',
                                    font=dict(color='#1565C0', size=16)
                                ),
                                xaxis=dict(title=col, gridcolor='#E3F2FD'),
                                yaxis=dict(title='Frequency', gridcolor='#E3F2FD'),
                                paper_bgcolor='#FFFFFF',
                                plot_bgcolor='#FFFFFF',
                                font=dict(color='#212121'),
                                showlegend=False,
                                height=400
                            )
                            
                            st.plotly_chart(fig, use_container_width=True)
                    
                    # Categorical distribution
                    if categorical_cols:
                        st.markdown("---")
                        st.markdown("#### Categorical Distribution")
                        st.markdown("**What this shows:** Comparison of different categories in your data. Longer bars indicate more frequent categories. This helps you see which groups are most common or rare.")
                        
                        for cat_col in categorical_cols[:3]:  # Show first 3 categorical columns
                            value_counts = df[cat_col].value_counts().head(10)
                            
                            fig = go.Figure(data=[go.Bar(
                                x=value_counts.values,
                                y=value_counts.index,
                                orientation='h',
                                marker_color='#2196F3',
                                marker_line_color='#1565C0',
                                marker_line_width=1.5
                            )])
                            
                            fig.update_layout(
                                title=dict(
                                    text=f'Top Values in {cat_col}',
                                    font=dict(color='#1565C0', size=16)
                                ),
                                xaxis=dict(title='Count', gridcolor='#E3F2FD'),
                                yaxis=dict(title=cat_col, gridcolor='#E3F2FD'),
                                paper_bgcolor='#FFFFFF',
                                plot_bgcolor='#FFFFFF',
                                font=dict(color='#212121'),
                                showlegend=False,
                                height=400
                            )
                            
                            st.plotly_chart(fig, use_container_width=True)
                    
                    # Correlation heatmap if multiple numeric columns
                    if len(numeric_cols) > 1:
                        st.markdown("---")
                        st.markdown("#### Correlation Matrix")
                        st.markdown("**What this shows:** How variables relate to each other. Darker blue means strong positive relationship (when one goes up, the other does too). White means no relationship. Numbers range from -1 to +1, with values closer to 1 or -1 indicating stronger relationships.")
                        
                        corr_matrix = df[numeric_cols].corr()
                        
                        fig = go.Figure(data=go.Heatmap(
                            z=corr_matrix.values,
                            x=corr_matrix.columns,
                            y=corr_matrix.columns,
                            colorscale=[[0, '#E3F2FD'], [0.5, '#FFFFFF'], [1, '#2196F3']],
                            zmid=0,
                            text=corr_matrix.values.round(2),
                            texttemplate='%{text}',
                            textfont={"size": 10},
                            colorbar=dict(title='Correlation')
                        ))
                        
                        fig.update_layout(
                            title=dict(
                                text='Feature Correlation Heatmap',
                                font=dict(color='#1565C0', size=18)
                            ),
                            paper_bgcolor='#FFFFFF',
                            plot_bgcolor='#FFFFFF',
                            font=dict(color='#212121'),
                            height=600
                        )
                        
                        st.plotly_chart(fig, use_container_width=True)
        
        # AR Visualisation Export
        st.markdown("---")
        st.markdown("### AR Visualisation Export")
        st.markdown("*Export 3D visualisations for AR viewing*")
        
        if df is not None and not df.empty:
            numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
            if len(numeric_cols) >= 2:
                if st.button("Generate AR-Ready 3D Model", type="primary"):
                    with st.spinner("Creating AR visualisation..."):
                        # Create 3D visualization
                        x_col = numeric_cols[0]
                        y_col = numeric_cols[1] if len(numeric_cols) > 1 else numeric_cols[0]
                        z_col = numeric_cols[2] if len(numeric_cols) > 2 else numeric_cols[0]
                        
                        fig = go.Figure(data=[go.Scatter3d(
                            x=df[x_col],
                            y=df[y_col],
                            z=df[z_col],
                            mode='markers',
                            marker=dict(
                                size=8,
                                color=df[z_col] if z_col in df.columns else '#2196F3',
                                colorscale=[[0, '#E3F2FD'], [0.5, '#1976D2'], [1, '#1565C0']],
                                line=dict(color='#1976D2', width=1),
                                showscale=True
                            ),
                            hovertemplate=f'{x_col}: %{{x}}<br>{y_col}: %{{y}}<br>{z_col}: %{{z}}<extra></extra>'
                        )])
                        
                        fig.update_layout(
                            title=dict(
                                text=f'AR-Ready 3D Visualisation',
                                font=dict(color='#1565C0', size=20)
                            ),
                            scene=dict(
                                xaxis=dict(title=x_col, backgroundcolor='#FFFFFF', gridcolor='#E3F2FD'),
                                yaxis=dict(title=y_col, backgroundcolor='#FFFFFF', gridcolor='#E3F2FD'),
                                zaxis=dict(title=z_col, backgroundcolor='#FFFFFF', gridcolor='#E3F2FD'),
                                camera=dict(
                                    eye=dict(x=1.5, y=1.5, z=1.5)
                                )
                            ),
                            paper_bgcolor='#FFFFFF',
                            font=dict(color='#212121'),
                            height=700
                        )
                        
                        st.plotly_chart(fig, use_container_width=True)
                        
                        # Export HTML for AR
                        html_buffer = BytesIO()
                        fig.write_html(html_buffer)
                        html_buffer.seek(0)
                        
                        st.download_button(
                            label="Download AR Visualisation (HTML)",
                            data=html_buffer,
                            file_name="ar_visualization.html",
                            mime="text/html",
                            help="Open this file on your phone/tablet with AR capabilities"
                        )
                        
                        st.info("**AR Viewing Tip:** Open the downloaded HTML file on a device with AR capabilities (smartphone/tablet). The 3D visualisation can be viewed in augmented reality mode.")
        else:
            st.info("AR visualisations are available for Excel, CSV, and JSON data with numeric columns.")
        
        # Download options
        st.markdown("---")
        st.markdown("### Download Options")
        
        col1, col2, col3 = st.columns(3)
        
        if df is not None and not df.empty:
            with col1:
                csv = df.to_csv(index=False).encode('utf-8')
                st.download_button(
                    label="Download Data (CSV)",
                    data=csv,
                    file_name="data_export.csv",
                    mime="text/csv"
                )
            
            with col2:
                # Basic statistics
                if len(df.select_dtypes(include=['number']).columns) > 0:
                    stats = df.describe().to_csv().encode('utf-8')
                    st.download_button(
                        label="Download Statistics",
                        data=stats,
                        file_name="data_statistics.csv",
                        mime="text/csv"
                    )
        
        with col3:
            # Download analysis report
            if 'analysis_text' in st.session_state:
                file_name = uploaded_file.name if uploaded_file else "Text Input"
                report = f"""MedInsight - Analysis Report
=====================================

File: {file_name}
Type: {file_type}
Generated: {pd.Timestamp.now().strftime('%Y-%m-%d %H:%M:%S')}

AI Analysis:
{st.session_state['analysis_text']}

Chat History:
"""
                for msg in st.session_state.get('chat_history', []):
                    report += f"\n{msg['role'].upper()}: {msg['content']}"
                
                st.download_button(
                    label="Download Full Report",
                    data=report.encode('utf-8'),
                    file_name="analysis_report.txt",
                    mime="text/plain"
                )
        
    except Exception as e:
        st.error(f"Error loading file: {str(e)}")
        st.info("Please make sure your file is a valid Excel (.xlsx, .xls) or CSV file")

else:
    # Instructions when no input provided
    st.markdown("### How to use:")
    st.markdown("""
    1. **Choose** between uploading a file or entering text directly
    2. **Upload** your file (Excel, CSV, PowerPoint, Text, JSON) OR **paste** text in the text box
    3. **Review** the data preview and details
    4. **Click** "Ask AI to Analyze" to get AI-powered insights
    5. **Ask questions** using the AI chatbot
    6. **Generate visualizations** including 3D AR graphs
    7. **Download** your results as needed
    """)
    
    st.markdown("### Supported formats:")
    st.markdown("""
    - **Data Files**: Excel (.xlsx, .xls), CSV
    - **PowerPoint**: .pptx (extracts tables and text)
    - **Text**: .txt, .md files
    - **JSON**: Structured data files
    - **Direct Text Input**: Paste any text for AI analysis
    """)
    
    st.markdown("### What you'll get:")
    st.markdown("""
    - **Quick overview** of your data or content structure
    - **AI-powered insights** about patterns and trends
    - **Key statistics** for numeric columns
    - **Clinical recommendations** for further analysis
    """)

# Footer
st.markdown("---")
st.markdown("*Powered by MedInsight AI*")

# Sidebar - Quick Navigation
with st.sidebar:
    st.markdown("### Quick Navigation")
    st.markdown("Jump to different sections:")
    
    if st.button("Chat (Top)", use_container_width=True):
        st.markdown('<a href="#ai-chat-assistant"></a>', unsafe_allow_html=True)
    
    if st.button("Upload Data", use_container_width=True):
        st.markdown('<a href="#choose-input-method"></a>', unsafe_allow_html=True)
    
    if st.button("AI Analysis", use_container_width=True):
        st.markdown('<a href="#ai-analysis"></a>', unsafe_allow_html=True)
    
    st.markdown("---")
    st.markdown("### Tips")
    st.markdown("""
    **Getting Started:**
    1. Use the chat above to ask questions
    2. Upload your data for analysis
    3. Get AI insights and visualizations
    
    **Best Practices:**
    - Start with the chat for quick questions
    - Upload data for detailed analysis
    - Use TTS for hands-free listening
    """)
